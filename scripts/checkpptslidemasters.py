import sys, os, csv, re, ipaddress
from lxml import etree as ET
from pptx.dml.color import RGBColor
# pip3 install python-pptx not pip3 install pptx
from cvss import CVSS2, CVSS3
import pandas as pd
import pprint
import itertools as it
import os, sys, re, shutil
#import datetime
import calendar, time
from pptx import Presentation
from pptx.util import Inches


def checkpptslidemasters():
    #print(1)
    wip_ppt=sys.argv[1]
    prs = Presentation(wip_ppt)
    slide_layout=int(sys.argv[2])
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))


# run this script from command line: python nexpose_netva_2_csv
if len (sys.argv) < 3 :
    #print("Usage: python "+os.path.basename(__file__) +" 'nexpose2csv' <inputfile>")
    print("Usage: python "+os.path.basename(__file__) +" wip_ppt layout_num")
    print("Exiting ...")
    sys.exit (1)
else:
    checkpptslidemasters()