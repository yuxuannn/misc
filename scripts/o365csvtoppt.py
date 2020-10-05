import sys, os, csv, re, ipaddress
from lxml import etree as ET
from pptx.dml.color import RGBColor
# pip3 install python-pptx not pip3 install pptx
from cvss import CVSS2, CVSS3
import pandas as pd
import pprint
import itertools as it
import os, sys, re, shutil
#import datetime
import calendar, time
from pptx import Presentation
from pptx.util import Inches

severity_colors={"CRITICAL": RGBColor(255, 0, 0),
                "HIGH": RGBColor(255, 102, 0),
                "MEDIUM": RGBColor(255, 192, 0),
                "LOW": RGBColor(46, 177, 53),
                "INFORMATION": RGBColor(0, 112, 192)
                }

severity_colors2={"Critical": RGBColor(255, 0, 0),
                "High": RGBColor(255, 102, 0),
                "Medium": RGBColor(255, 192, 0),
                "Low": RGBColor(46, 177, 53),
                "Information": RGBColor(2, 112, 192)
                }


def convertCVSS2toCVSS3(cvss2):
    #cvss2='AV:L/AC:M/Au:S/C:C/I:C/A:P'
    cvss2arr=cvss2.split('/')
    cvss3='AV:P/AC:L/PR:N/UI:N/S:U/C:N/I:N/A:N'
    #print(cvss2)
    #print(cvss3)
    #print(cvss2arr)
    
    cvss3arr=cvss3.split('/')
    cvss3arr[0]=cvss2arr[0]
    if cvss2arr[1]=='AC:L':
        cvss3arr[1]=cvss2arr[1]
    else:
        cvss3arr[1]='AC:H'

    if cvss2arr[2]=='Au:N':
        cvss3arr[2]='PR:N'
    else:
        cvss3arr[2]='PR:L'


    for num in range(3,6):
        #print(num)
        #print(cvss2arr[num], cvss3arr[num+2])
        (metric, value)=cvss2arr[num].split(':')
        #print(value)
        if value=='N':
            cvss3arr[num+2]=metric+':'+'N'
        elif value=='P':
            cvss3arr[num+2]=metric+':'+'L'
        elif value=='C':
            cvss3arr[num+2]=metric+':'+'H'    



    #print(cvss3arr)
    cvss3='CVSS:3.0/'+'/'.join(cvss3arr)
    #print(cvss3)
    return cvss3

def mapCVSS3toScore(cvss3Vector):
    c = CVSS3(cvss3Vector)
    return c.scores()[0]

# parse complicated nodes with different nested tags
def scrapenode(elem, tag):
    #print(elem.attrib)
    fullvuln=[]
    
    # iterate through all child elements recursively?
    for element in elem.getiterator():
        
        #if element.tag == 'solution':
        ''' old
        if element.tag == tag:
            continue
        else:
            tmp=[''+key+ ' : '+element.attrib.get(key) for key in element.attrib.keys()]
            fullvuln.append(str(tmp))
        
        '''
        if element.tag == 'URLLink':
            fullvuln.append(element.attrib.get('LinkURL'))
        
        if element.text:
            
            text=re.sub(r'\n+', '\n', element.text).strip()
            text=re.sub(r'\r+', '\r', text).strip()
            # remove duplicated space
            text=re.sub(' +', ' ', text).strip()
            #print(element.text)
            fullvuln.append(text)

    fullvuln='\n'.join(fullvuln).replace('\r', '').replace('\t','').replace('[]','')
    # /replace-multiple-newlines-with-single-newlines-during-reading-file
    fullvuln = re.sub(r'\n+', '\n', fullvuln).strip()
    
    return fullvuln

    
def mapassettype(ipaddr):
    try:
        addr = ipaddress.IPv4Address(ipaddr)
    except ValueError:
        #raise # not an IP address
        return 'UNKNOWN'
    if addr.is_private:
        return 'INT' # is a private address
    else:
        return 'EXT'
    '''
    if ipaddr in ['61.8.245.52','122.11.168.160','61.8.245.50','61.8.245.53','40.90.176.90','202.79.195.76','202.79.195.81','202.79.195.82']:
        return 'EXT'
    elif ipaddr in ['10.90.1.39','10.90.1.37','10.90.1.31','10.90.1.65','10.90.1.2','10.90.1.28','10.90.15.5','192.168.90.6']:
        return 'INT'
    else:
        return 'UNKNOWN'
    '''


def CVSSSeverity(score):
    # convert back to float 
    score=float(score)
    if score >= 9.0 and score <= 10.0:
        severity='CRITICAL'
    elif score >= 7.0 and score <= 8.9:
        severity='HIGH'
    elif score >= 4.0 and score <= 6.9:
        severity='MEDIUM'        
    elif score >= 0.1 and score <= 3.9:
        severity='LOW'
    elif score == 0:
        severity='NONE'
    else:
        raise ValueError()
    return severity

def CVSSSeverity2(score):
    # convert back to float 
    score=float(score)
    if score >= 9.0 and score <= 10.0:
        severity='Critical'
    elif score >= 7.0 and score <= 8.9:
        severity='High'
    elif score >= 4.0 and score <= 6.9:
        severity='Medium'        
    elif score >= 0.1 and score <= 3.9:
        severity='Low'
    elif score == 0:
        severity='Information'
    else:
        raise ValueError()
    return severity

# Map HTTP/HTTPS to HTTP(S) and combine by Issue Name (aggregate Issue Evidence too)
def mapsvc(x):
    if x=='HTTP':
        return 'HTTP(S)'
    elif x=='HTTPS':
        return 'HTTP(S)'
    else:
        return x

# Add Service Name to Title as prefix if not present within title
def addsvcname2title(data2):
    #print (data2.apply(lambda x: x['Service Name'] in x['Description'] and x['Service Name'] != "'", axis=1).head(5))
    searchindex=data2.index[data2.apply(lambda x: x['Service Name'] not in x['title'] and x['Service Name'] != "'", axis=1)].tolist()
    print(searchindex)
    #tmp= data2.loc[searchindex][['Service Name','Description']]
    #data2tmp=data2.loc[searchindex].apply(lambda x: x['title']='(' +x['Service Name']+ ')'+x['title'], axis=1)
    #data2.loc[searchindex].apply(lambda x: print(x['title']), axis=1)
    #data2.loc[searchindex]['title']='(' +data2.loc[searchindex]['Service Name']+ ')'+data2.loc[searchindex]['title']
    data2tmp=data2.loc[searchindex]
    #data2.loc[searchindex]=data2tmp.assign(title='(' +data2.loc[searchindex]['Service Name']+ ') '+data2.loc[searchindex]['title'])
    data2.loc[searchindex]=data2tmp.assign(title=data2.loc[searchindex]['Service Name']+ ' - '+data2.loc[searchindex]['title'])
    #return data2tmp



# Compress rows according to asset type (if defined, else default REF)
def merge_rows4(df, searchindex, newcol, counter):
    
    if searchindex:
        #print(searchindex)
        df1 = df.loc[searchindex]
        df1['key'] = 'a'
        #df2=df1.groupby(['key']).transform(lambda x : '**\n'.join(x))
        #print(df1[['title', 'cvssScore']])
        #print(df1['cvssScore'].max())
        # simplify risk
        #for risk in ['Critical', 'High', 'Medium', 'Low']:
            #df2['Risk'][df2['Risk'].str.contains(risk, case=False)]=risk

        #df2['Host:port']=df2['Host:port'].map(lambda x: simplifyhostport(x))

        df1['REF']=newcol
        #df1['REF-ID']='REF-'+str(counter)
        df1['REF-ID']=df1['AssetType']+'-'+str(counter)
        print(df1['cvssScore'].max())
        df1['Overall CVSS']=df1['cvssScore'].max()
        
        df.loc[searchindex]=df1
        #print(df2)
        #df.loc[searchindex]=df2        
        #return df

def backupFile(wip_ppt):
    inputFile=wip_ppt
    (baseFolder, basename)=os.path.split(inputFile)
    print(baseFolder, basename)
    (filename, fileext)=os.path.splitext(basename)
    print(filename, fileext)

    #outputFile=os.path.join(baseFolder, filename+'_output'+fileext)

    # need to cast to string, if not get this error "TypeError: can only concatenate str (not "int") to str"
    #ts = str(datetime.datetime.now().timestamp()) # 1581597259.2283
    ts = str(calendar.timegm(time.gmtime())) # 1581597259
    outputFile=os.path.join(baseFolder, filename+'_'+ts+fileext)
    print(outputFile)
    # will override existing file
    shutil.copyfile(inputFile, outputFile) 

def checkpptslidemasters():
    print(1)
    wip_ppt=sys.argv[2]
    prs = Presentation(wip_ppt)
    slide_layout=int(sys.argv[3])
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))


def addWebIssueSlide(info, slide_layout):
            
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
    # 1 Ref Title
    slide.shapes.title.text = "REF-01"
    # 2 Title/Issue Name
    slide.placeholders[11].text = title # Title
    # 3 Affected Hosts
    slide.placeholders[20].text = affected_hosts # Host
    # 4 Risk Rating
    slide.placeholders[23].fill.solid()
    slide.placeholders[23].transparency = 0
    slide.placeholders[23].fill.fore_color.rgb = severity_colors2[risk_rating]
    slide.placeholders[23].text = risk_rating # Risk Rating
    # Font White Color
    slide.placeholders[23].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    # 5 CVSS Base Score
    slide.placeholders[24].text = str(cvss_base) +' [' +cvss_vector +']'# CVSS Base Score
    # 6 OWASP Category
    slide.placeholders[28].text = owaspcat # OWASP Category
    # 7 Description and Risks
    slide.placeholders[21].text = description_impact # Description & Risks
    # 8 Remediation & Further Info
    slide.placeholders[22].text = remediation # Remediation & Further Info
    

def addOIssueSlide(info, prs, slide_layout):
    # info = [Ref, Title, Affected Hosts, Risk Rating, CVSS Base Score, OWASP Category, Description & Risks, Remediation and Further Info]
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
    # 1 Ref Title
    slide.shapes.title.text = info[0]
    # 2 Title/Issue Name
    slide.placeholders[11].text = info[1]
    # 3 Affected Hosts
    slide.placeholders[20].text = info[2]
    # 4 Risk Rating
    slide.placeholders[23].fill.solid()
    slide.placeholders[23].transparency = 0
    slide.placeholders[23].fill.fore_color.rgb = severity_colors2[info[3]]
    slide.placeholders[23].text = info[3]
    # Font White Color
    slide.placeholders[23].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    # 5 CVSS Base Score
    slide.placeholders[24].text = info[4]
    # 6 OWASP Category
    slide.placeholders[28].text = info[5]
    # 7 Description and Risks
    slide.placeholders[21].text = info[6]
    # 8 Remediation & Further Info
    slide.placeholders[22].text = info[7]



            

def addOTestOutputSlide(info, prs, slide_layout):
    #print(123)
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
    # 1 Ref Title
    slide.shapes.title.text = info[0]
    slide.placeholders[11].text = info[1]

def o365csvtoppt():
    #print(1)
    #csvinput="shareableasets-webvuln1.csv"
    csvinput=sys.argv[1]
    #csvinput="VulnKB2020.csv"
    #wip_ppt='template_occ.pptx'
    wip_ppt=sys.argv[2]

    backupFile(wip_ppt)

    #prs = Presentation('template_occ.pptx')

    prs = Presentation(wip_ppt)
    # unused slides will be gone, hence numbering may change per customised template
    #vuln_web_slide_layout=5
    vuln_web_slide_layout=4
    issue_slide_layout=6
    testoutput_slide_layout=5

    
    with open(csvinput, errors='ignore') as f:
        reader = csv.DictReader(f, delimiter=',')
        #print('type of reader: ', type(reader))
        

        for row in reader:
            #reported=row['Reported']
            
            #if reported=='y':
                #continue
            
            
            #print('type of row: ', type(row))
            #print(row)
            #print(row['Risk Rating'])  # Access by column header instead of column number
            #print(row['Issue Title'])
            

            #risk_rating=row['Risk']
            # info = [Ref, Title, Affected Hosts, Risk Rating, CVSS Base Score, OWASP Category, Description & Risks, Remediation and Further Info]
            info=[0,1,2,3,4,5,6,7]
            info[0]=row['REF']
            
            
            info[1]=row['Title']
            info[2]=row['Affected Product']
            info[3]=row['Severity']
            info[4]=row['Category']
            info[5]="Pass/Fail/NA"
            #info[6]=row['Description Notes']+' ' +row['Impact']
            # add line break - works
            info[6]=row['Description']+"\n" +row['Impact']
            info[7]=row['Recommendation']

            #print(info)
            addIssueSlide(info, prs, issue_slide_layout)

            #print(row['CMD?'])
            #print(row['Top-Level Function'])
            if row['UI Navigation Check']=='':
                print('UI Navigation Check is Blank')
            else:
                #print('Top-Level Function is not Blank', row['UI Navigation Check'])
                info2=[0,1]
                info2[0]=row['REF']
                #info2[0]=issueref
                info2[1]="Screenshot: "+ row['UI Navigation Check']
                addTestOutputSlide(info2, prs, testoutput_slide_layout)

            if row['CMD?']=='Y':
                
                info2=[0,1]
                info2[0]=row['REF']
                #info2[0]=issueref
                #info2[1]="Powershell Command: "+row['Command']+"\n"+"Command Output: "
                info2[1]="Powershell Command: "+row['Command']
                addTestOutputSlide(info2, prs, testoutput_slide_layout)
            #addO365IssueSlide(info, issue_slide_layout)


            #addO365IssueSlide()
            #addO365TestOutputSlide()
            

            
            

    #prs.save('template8.pptx')

    prs.save(wip_ppt)


# run this script from command line: python nexpose_netva_2_csv
if len (sys.argv) < 2 :
    #print("Usage: python "+os.path.basename(__file__) +" 'nexpose2csv' <inputfile>")
    print("Usage: "+os.path.basename(__file__) +" 'csvtoppt' csvinput wip_ppt")
    print("Exiting ...")
    sys.exit (1)
else:
    o365csvtoppt()