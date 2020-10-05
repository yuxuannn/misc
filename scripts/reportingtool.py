#! /usr/bin/python3

import itertools as it
import os, sys, re, shutil
#import datetime
import calendar, time
from pptx import Presentation
from pptx.util import Inches
import csv
from pptx.dml.color import RGBColor
import pandas
import pandas as pd
import pprint
from lxml import etree as ET
from reportingtoolhelperfuncs import *
#from nexpose_netva_2_csv import *

if len (sys.argv) < 2 :
    #print("Usage: python "+os.path.basename(__file__) +" 'nexpose2csv' <inputfile>")
    print("Usage: "+os.path.basename(__file__) +" 'nexpose2csv' <inputfile>")
    print("Usage: "+os.path.basename(__file__) +" 'nexposepolicy2csv' xccdf.xml")
    print("Usage: "+os.path.basename(__file__) +" 'csvtoppt' csvinput wip_ppt")
    print("Exiting ...")
    sys.exit (1)



'''

#print("Running Python Script processnmap.py")
#filename='./nmap-output.txt'
inputFile=sys.argv[1]
print(inputFile)

(baseFolder, basename)=os.path.split(inputFile)
print(baseFolder, basename)
(filename, fileext)=os.path.splitext(basename)
print(filename, fileext)

#outputFile=os.path.join(baseFolder, filename+'_output'+fileext)

# need to cast to string, if not get this error "TypeError: can only concatenate str (not "int") to str"
#ts = str(datetime.datetime.now().timestamp()) # 1581597259.2283
ts = str(calendar.timegm(time.gmtime())) # 1581597259
outputFile=os.path.join(baseFolder, filename+'_'+ts+fileext)
print(outputFile)
# will override existing file
shutil.copyfile(inputFile, outputFile) 



'''
def backupFile(wip_ppt):
    inputFile=wip_ppt
    (baseFolder, basename)=os.path.split(inputFile)
    print(baseFolder, basename)
    (filename, fileext)=os.path.splitext(basename)
    print(filename, fileext)

    #outputFile=os.path.join(baseFolder, filename+'_output'+fileext)

    # need to cast to string, if not get this error "TypeError: can only concatenate str (not "int") to str"
    #ts = str(datetime.datetime.now().timestamp()) # 1581597259.2283
    ts = str(calendar.timegm(time.gmtime())) # 1581597259
    outputFile=os.path.join(baseFolder, filename+'_'+ts+fileext)
    print(outputFile)
    # will override existing file
    shutil.copyfile(inputFile, outputFile) 

def checkpptslidemasters():
    print(1)
    wip_ppt=sys.argv[2]
    prs = Presentation(wip_ppt)
    slide_layout=int(sys.argv[3])
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))

def csvtoppt():
    print(1)
    #csvinput="shareableasets-webvuln1.csv"
    csvinput=sys.argv[2]
    #csvinput="VulnKB2020.csv"
    #wip_ppt='template_occ.pptx'
    wip_ppt=sys.argv[3]

    backupFile(wip_ppt)

    #prs = Presentation('template_occ.pptx')

    prs = Presentation(wip_ppt)
    # unused slides will be gone, hence numbering may change per customised template
    #vuln_web_slide_layout=5
    vuln_web_slide_layout=4

    import csv
    with open(csvinput) as f:
        reader = csv.DictReader(f, delimiter=',')
        print('type of reader: ', type(reader))
        

        for row in reader:
            reported=row['Reported']
            
            if reported=='y':
                continue
            
            
            #print('type of row: ', type(row))
            #print(row)
            #print(row['Risk Rating'])  # Access by column header instead of column number
            #print(row['Issue Title'])
            

            #risk_rating=row['Risk']

            title=row['Issue Name']
            #description_impact=row['Description and Risks']
            description_impact=row['Description']+' ' +row['Impact']
            remediation=row['Remediation and Further Information']
            #cvss_base=row['CVSS3 Base Score']

            cvss_vector=row['CVSS3 Vector']
            cvss_base=mapCVSS3toScore(cvss_vector)
            risk_rating=CVSSSeverity2(cvss_base)

            owaspcat=row['OWASP Category']
            affected_hosts=row['Affected Hosts']
            
            slide = prs.slides.add_slide(prs.slide_layouts[vuln_web_slide_layout])
            slide.shapes.title.text = "REF-01"
            slide.placeholders[11].text = title # Title
            slide.placeholders[20].text = affected_hosts # Host
            slide.placeholders[21].text = description_impact # Description & Risks
            slide.placeholders[22].text = remediation # Remediation & Further Info
            
            slide.placeholders[24].text = str(cvss_base) +' [' +cvss_vector +']'# CVSS Base Score
            slide.placeholders[28].text = owaspcat # OWASP Category

            slide.placeholders[23].fill.solid()
            slide.placeholders[23].transparency = 0
            slide.placeholders[23].fill.fore_color.rgb = severity_colors2[risk_rating]
            slide.placeholders[23].text = risk_rating # Risk Rating
            # Font White Color
            slide.placeholders[23].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            

    #prs.save('template8.pptx')

    prs.save(wip_ppt)


def nexposepolicy2csv():
    print(3)
    #inputfile='xccdf.xml'
    inputFile=sys.argv[2]
    (baseFolder, basename)=os.path.split(inputFile)
    print(baseFolder, basename)
    (filename, fileext)=os.path.splitext(basename)
    print(filename, fileext)

    outputFile=os.path.join(baseFolder, 'policycheckresultswithdetails_'+filename+'.csv')

    #inputfile='st telemedia - internal servers - xml export 2.0 report.xml'
    #import xml.etree.ElementTree as ET
    # use lxml instead because got getparent function

    tree=ET.parse(inputFile)
    root=tree.getroot()

    df1=extractnexposepolicyvulndefns(root)
    df2=extractnexposepolicyfindings(root)
    df3=pd.merge(df1,df2, left_on='id', right_on='id')
    df3.to_csv(outputFile)
    print('Output File: '+outputFile)
    



















options = { 'csvtoppt': csvtoppt,
            #'nexpose2csv': nexpose2csv,
            'nexposepolicy2csv': nexposepolicy2csv,
	        'checkpptslidemasters': checkpptslidemasters
		

}

option=sys.argv[1]
options[option]()

